"""Generate the ParkSight pitch deck (.pptx) — dark, judge-ready, embeds figures."""
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))
from parksight import config as C  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402

BG = RGBColor(0x0B, 0x0F, 0x1A)
PANEL = RGBColor(0x16, 0x1D, 0x2E)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
GREY = RGBColor(0x94, 0xA3, 0xB8)
INDIGO = RGBColor(0x6B, 0x6E, 0xF1)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xFA, 0xCC, 0x15)

SW, SH = Inches(13.333), Inches(7.5)
meta = json.loads((C.PROCESSED / "meta.json").read_text())
mx = json.loads((C.MODELS.parent / "metrics.json").read_text())

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def txt(s, x, y, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        italic=False, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb


def bar(s, color=INDIGO, y=Inches(1.18), h=Pt(4)):
    b = s.shapes.add_shape(1, Inches(0.6), y, Inches(2.2), h)
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()


def bullets(s, x, y, w, h, items, size=18, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "•  " + it
        r.font.size = Pt(size); r.font.color.rgb = WHITE if i % 1 == 0 else GREY
        r.font.name = "Segoe UI"
    return tb


def pic(s, path, x, y, w=None, h=None):
    if Path(path).exists():
        s.shapes.add_picture(str(path), x, y, width=w, height=h)


def header(s, kicker, title, color=INDIGO):
    txt(s, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4), kicker, 14, color, bold=True)
    txt(s, Inches(0.6), Inches(0.75), Inches(12.1), Inches(0.9), title, 30, WHITE, bold=True)
    bar(s, color)


# 1 — Title
s = slide()
txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.2), "🅿️  ParkSight", 60, WHITE, bold=True)
txt(s, Inches(0.95), Inches(3.5), Inches(11.5), Inches(0.8),
    "AI-Driven Parking-Congestion Intelligence for Targeted Enforcement", 24, CYAN, bold=True)
txt(s, Inches(0.95), Inches(4.4), Inches(11.5), Inches(0.8),
    "From reactive patrols to a predictive command center —\nbuilt on "
    f"{meta['total_violations']:,} real Bengaluru parking violations.", 18, GREY)
txt(s, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.5),
    "Flipkart Gridlock 2.0 · Round 2 · Theme 1", 14, GREY, italic=True)

# 2 — Problem
s = slide(); header(s, "THE CHALLENGE", "Parking chokes the city — and nobody can see it")
bullets(s, Inches(0.6), Inches(1.7), Inches(7.2), Inches(4.5), [
    "On-street & spillover parking near markets, metro and events blocks carriageways and junctions.",
    "Enforcement today is patrol-based and reactive — officers chase, they don't pre-empt.",
    "No heatmap of violations vs. congestion impact exists.",
    "Authorities can't prioritise which zones to enforce, or when.",
], 19, gap=14)
txt(s, Inches(8.1), Inches(1.8), Inches(4.6), Inches(4),
    "“Where does illegal parking\nactually hurt traffic flow —\nand where do we send teams\ntomorrow?”",
    22, AMBER, italic=True, bold=True)

# 3 — Insight
s = slide(); header(s, "THE DISCOVERY", "The Evening Enforcement Blind-Spot", RED)
pic(s, C.ASSETS / "fig_blindspot.png", Inches(0.6), Inches(1.7), w=Inches(7.4))
bullets(s, Inches(8.2), Inches(1.9), Inches(4.7), Inches(4.5), [
    "After converting timestamps UTC→IST, records collapse during 15:00–24:00.",
    "That's exactly when commercial congestion peaks.",
    "Few tickets ≠ few violations → it's low enforcement VISIBILITY.",
    "This IS the 'poor visibility' the brief names — ParkSight closes it.",
], 17, gap=12)

# 4 — Dataset
s = slide(); header(s, "THE DATA", "298,450 violations — a near-perfect fit", CYAN)
stats = [(f"{meta['total_violations']:,}", "violations"),
         ("100%", "geo-located"),
         (f"{meta['n_stations']}", "police stations"),
         (f"{meta['n_junctions']}", "junctions"),
         (f"{meta['repeat_share']}%", "from repeat offenders"),
         (f"{meta['severe_share']}%", "carriageway-blocking")]
for i, (v, l) in enumerate(stats):
    x = Inches(0.6 + (i % 3) * 4.15); y = Inches(1.9 + (i // 3) * 1.9)
    card = s.shapes.add_shape(1, x, y, Inches(3.9), Inches(1.6))
    card.fill.solid(); card.fill.fore_color.rgb = PANEL; card.line.color.rgb = INDIGO
    card.shadow.inherit = False
    txt(s, x, y + Inches(0.18), Inches(3.9), Inches(0.8), v, 34, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, y + Inches(1.0), Inches(3.9), Inches(0.5), l, 15, GREY, align=PP_ALIGN.CENTER)
txt(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.8),
    f"Nov 2023 → Apr 2024 · dominated by Wrong/No-Parking & Main-Road parking · "
    f"{meta.get('n_devices',0):,} enforcement devices.", 15, GREY)

# 5 — Solution overview
s = slide(); header(s, "THE PLATFORM", "Eight modules a command center would actually use")
mods = [("🧭 Executive Dashboard", "KPIs, impact, the blind-spot"),
        ("🗺️ AI Hotspot Map", "H3 heat, density, zones, drill-down"),
        ("📊 PCIS Engine", "explainable impact score, re-weightable"),
        ("🔮 Forecast", "next-7-day per station + confidence"),
        ("🎯 Prioritize", "ranked plan + PDF briefing"),
        ("🧪 Simulator", "What-If deployment impact"),
        ("🤖 AI Copilot", "plain-English, grounded answers"),
        ("🚨 Offenders", "15% of vehicles → 34% of violations")]
for i, (t, d) in enumerate(mods):
    x = Inches(0.6 + (i % 4) * 3.12); y = Inches(2.0 + (i // 4) * 2.3)
    card = s.shapes.add_shape(1, x, y, Inches(2.9), Inches(2.0))
    card.fill.solid(); card.fill.fore_color.rgb = PANEL; card.line.color.rgb = INDIGO
    card.shadow.inherit = False
    txt(s, x + Inches(0.12), y + Inches(0.15), Inches(2.7), Inches(0.9), t, 15, WHITE, bold=True)
    txt(s, x + Inches(0.12), y + Inches(1.0), Inches(2.7), Inches(0.9), d, 12.5, GREY)

# 6 — PCIS
s = slide(); header(s, "THE MOAT", "PCIS — quantifying impact without faking traffic data", AMBER)
txt(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.7),
    "PCIS = 100 × (0.30·Volume + 0.20·Severity + 0.20·Location + 0.20·PeakOverlap + 0.10·Trend)",
    20, CYAN, bold=True)
bullets(s, Inches(0.6), Inches(2.7), Inches(12), Inches(3.5), [
    "Volume — log-scaled violation count (heavy-tailed).",
    "Severity — carriageway-blocking weight: main-road / double / footpath > no-parking.",
    "Location — junction & commercial-keyword criticality (a junction block cascades).",
    "Peak-overlap — share in IST peak windows 08–11 & 17–21 (a 6 PM block hurts more than 3 AM).",
    "Trend — month-over-month growth, so rising hotspots are caught early.",
    "Weights are TUNABLE live in the app — a transparent policy lever, not a black box.",
], 18, gap=12)

# 7 — Hotspots
s = slide(); header(s, "SEE IT", "Hotspots ranked by real impact, not raw counts")
_map = C.ASSETS / "screens" / "02_hotspot_map.png"
pic(s, _map if _map.exists() else C.ASSETS / "fig_top_zones.png",
    Inches(0.6), Inches(1.7), w=Inches(7.6))
bullets(s, Inches(8.4), Inches(2.0), Inches(4.5), Inches(4), [
    "H3 hexagons + DBSCAN zones on a dark deck.gl map.",
    "Colour = PCIS, so a main-road evening hotspot beats a quiet residential one.",
    "Click any zone → component breakdown + recommended units & time window.",
], 17, gap=14)

# 8 — Forecast
s = slide(); header(s, "PREDICT", "Tomorrow's hotspots, with confidence", CYAN)
pic(s, C.ASSETS / "fig_forecast.png", Inches(0.6), Inches(1.7), w=Inches(7.6))
bullets(s, Inches(8.4), Inches(2.0), Inches(4.6), Inches(4), [
    f"LightGBM, 7-day horizon, p10–p90 confidence bands.",
    f"MAE {mx['mae_model']} vs baseline {mx['mae_baseline']} — {mx['improvement_pct']}% better.",
    "Benchmarked against climatology so the ML's value is explicit.",
    "Feature importance = explainability, not a black box.",
], 17, gap=14)

# 9 — Act
s = slide(); header(s, "ACT", "Prioritise → Simulate → Brief", RED)
bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), [
    "Prioritize: ranked High/Med/Low per station & junction, each with a plain-English reason.",
    "Enforcement-Gap targets where predicted impact is high but current presence is low — the blind spot.",
    "Smart Simulator: pick zones, set an enforcement increase & deterrence elasticity → see violations prevented and % of citywide impact covered, live.",
    "One-click Daily Deployment Briefing PDF — a station officer can act on it this shift.",
], 19, gap=16)

# 10 — Copilot
s = slide(); header(s, "ASK", "An AI Copilot, grounded in the data")
bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(3), [
    "“Show the top 10 illegal parking zones.”",
    "“Which hotspot causes the highest congestion?”",
    "“Where should enforcement teams deploy tomorrow evening?”",
    "“What happens if we increase enforcement by 20%?”",
], 20, gap=14)
txt(s, Inches(0.6), Inches(5.2), Inches(12), Inches(1.5),
    "Claude parses the question; our analytics compute the answer — every number is grounded in the "
    "dataset, never hallucinated. A deterministic fallback keeps the demo alive with no API key.",
    17, GREY, italic=True)

# 11 — Architecture
s = slide(); header(s, "HOW IT'S BUILT", "Precompute → tiny artifacts → instant app", CYAN)
pic(s, C.ASSETS / "architecture.png", Inches(0.9), Inches(1.7), w=Inches(11.5))

# 12 — Close
s = slide(); header(s, "IMPACT & HONESTY", "Deployable, defensible, and honest")
bullets(s, Inches(0.6), Inches(1.7), Inches(7.4), Inches(4.6), [
    "A working, deployed product the Bengaluru Traffic Police could pilot now.",
    "PCIS is an explicit proxy — no traffic-flow column exists, and we say so.",
    "No closure/action timestamps → we don't overclaim response-time analytics.",
    "Roadmap: live SCITA/ITMS feed, OSM road-graph spillover, cell-level forecasting, patrol optimiser.",
], 18, gap=14)
txt(s, Inches(8.2), Inches(1.9), Inches(4.6), Inches(4),
    "ParkSight\nsees the blind spot,\nscores the impact,\nand deploys the plan.", 24, AMBER, bold=True)

out = C.REPORTS_OUT / "ParkSight_Pitch_Deck.pptx"
prs.save(str(out))
print(f"wrote {out} ({out.stat().st_size/1024:.0f} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
